import os
import uuid 
import time
import openai
import ast
import re
import logging
import io
from io import BytesIO
from dotenv import load_dotenv
from urllib.parse import urlparse
from pydantic import BaseModel, Field
from langchain_core.prompts import PromptTemplate
import asyncio
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict
from datetime import datetime
from flask_cors import CORS
import numpy as np
import pandas as pd
import requests 
import docx2txt2
import xlrd
import json
import subprocess
import concurrent.futures
from functools import partial
from sklearn.metrics.pairwise import cosine_similarity
from openpyxl.styles import Alignment
from flask import Flask, request, jsonify, session
from langchain.agents import Tool, create_openai_functions_agent, AgentExecutor
from langchain.schema import Document,BaseRetriever
from langchain_openai import AzureChatOpenAI, AzureOpenAIEmbeddings
from langchain.chains import RetrievalQA
from langchain_community.callbacks import get_openai_callback
from langchain_community.retrievers import BM25Retriever
from langchain.retrievers import EnsembleRetriever
from langchain_community.vectorstores import FAISS
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_text_splitters import MarkdownHeaderTextSplitter
from pptx import Presentation
from azure.storage.blob import BlobServiceClient
from utilities.Document_Processing import BlobStorageProcessor, DocumentIntelligenceLoader
from utilities.ABB_entity_extraction import extract_entities_from_blob
from utilities.Excel_Formatting import create_excel_with_formatting_local
from utilities.prompts import prompts
from utilities.facilites import read_xls_to_prompts, extract_and_parse_json, process_single_file
from utilities.materials import ExcelMatcher, read_excel_preserve_format
from utilities.facilites_xlsx import facilitesMatcher, fix_duplicate_columns
from tempfile import NamedTemporaryFile
from langchain.globals import set_debug
from datetime import datetime

set_debug(True)

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__,template_folder='template',static_folder='./frontend/build',static_url_path='/')
app.json.sort_keys = False
CORS(app)

# Configuration for Azure OpenAI
openai.api_type = os.getenv("AZURE_OPENAI_TYPE")
openai.api_base = os.getenv('AZURE_OPENAI_ENDPOINT')
openai.api_key = os.getenv('AZURE_OPENAI_API_KEY')
openai.api_version = os.getenv('AZURE_OPENAI_API_VERSION')

# Azure configuration
LLM_MODEL = os.getenv('AZURE_OPENAI_GPT4_DEPLOYMENT_NAME')
AZURE_BLOB_STORAGE_CONTAINER_NAME = os.getenv('AZURE_BLOB_STORAGE_CONTAINER_NAME')
AZURE_BLOB_STORAGE_ACCOUNT_NAME = os.getenv('AZURE_BLOB_STORAGE_ACCOUNT_NAME')
AZURE_BLOB_STORAGE_ACCOUNT_KEY = os.getenv('AZURE_BLOB_STORAGE_ACCOUNT_KEY')
AZURE_BLOB_STORAGE_CONTAINER_URL = os.getenv('AZURE_BLOB_STORAGE_CONTAINER_URL')
AZURE_VECTOR_STORAGE_CONTAINER_NAME = os.getenv('AZURE_VECTOR_STORAGE_CONTAINER_NAME')
AZURE_BLOB_STORAGE_CONNECTION_STRING = os.getenv('AZURE_BLOB_STORAGE_CONNECTION_STRING')
AZURE_PRELOADED_CONTAINER_NAME = os.getenv('AZURE_PRELOADED_CONTAINER_NAME')
AZURE_DOWNLOAD_STORAGE_CONTAINER_NAME = os.getenv('AZURE_DOWNLOAD_STORAGE_CONTAINER_NAME')
AZURE_STORAGE_CONNECTION_STRING_DOWNLOAD_EXCEL = os.getenv('AZURE_STORAGE_CONNECTION_STRING_DOWNLOAD_EXCEL')
AZURE_DOWNLOAD_EXCEL_STORAGE_CONTAINER_NAME = os.getenv('AZURE_DOWNLOAD_EXCEL_STORAGE_CONTAINER_NAME')
AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT_NAME = os.getenv('AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT_NAME')
AZURE_MODEL = os.getenv('AZURE_OPENAI_MODEL')
AZURE_OPENAI_EMBEDDINGS_VERSION = os.getenv('AZURE_OPENAI_EMBEDDINGS_VERSION')

# Initialize services
llm = AzureChatOpenAI(
    azure_deployment=LLM_MODEL,
    openai_api_key=openai.api_key,
    openai_api_version=openai.api_version,
    azure_endpoint=openai.api_base
)

blobstorageprocessor = BlobStorageProcessor()
documentintelligence = DocumentIntelligenceLoader()
embeddings = AzureOpenAIEmbeddings(
    azure_deployment=AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT_NAME,
    openai_api_version=openai.api_version,
)

GLOBAL_PROMPTS = None

class Information(BaseModel):
    Entity: str = Field(title="Summary", description="Extract the relevant information from the document. If not explicitly provided, do not guess.")

class DocumentInput(BaseModel):
    question: str = Field()

class EmptyRetriever(BaseRetriever):
    def _get_relevant_documents(self, query: str):
        return []

    async def _aget_relevant_documents(self, query: str):
        return []

matcher = ExcelMatcher()

class OrderedEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, OrderedDict):
            return {k: self.default(v) for k, v in obj.items()}
        if isinstance(obj, np.integer):
            return int(obj)
        if isinstance(obj, np.floating):
            return float(obj)
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        return super().default(obj)
        
def id_generator():
    return str(uuid.uuid4())

def convert_to_dict(text):
    start = text.index('{')
    end = text.rfind('}') + 1
    try:
        return ast.literal_eval(text[start:end])
    except SyntaxError as e:
        return {"error": f'Error parsing LLM response: {e}'}

def string_to_dict(string_list):
    """
    Convert a string representation of a list into a dictionary with "additions" as the key
    and the list of strings as the value, formatted according to the specified output format.
    
    Args:
    string_list (str): String representation of a list to be converted into a dictionary.
    
    Returns:
    dict: Dictionary with "additions" as the key and the list of formatted strings as the value.
    """
    # Convert the string representation of the list into an actual list
    actual_list = ast.literal_eval(string_list)
    
    # Split the string into individual summaries
    summaries = actual_list[0].split(",\n")
    
    # Convert the list into a dictionary with "additions" as the key
    return {"additions": summaries}

@app.route('/')
def index():
    logger.info('Index route accessed.')
    return app.send_static_file('index.html')

blob_service_client = BlobServiceClient.from_connection_string(AZURE_BLOB_STORAGE_CONNECTION_STRING)
blob_list=[]

@app.route('/upload', methods=['POST'])
def upload_files():
    try:
        # Check if the request contains 'index_name' as a boolean indicator
        index_name = request.form.get('index_name') is not None
        supplier_gtc = request.files.get('supplier_gtc')
        abb_gtc = request.files.get('abb_gtc') if not index_name else None

        if not supplier_gtc:
            logger.error('Supplier GTC file must be uploaded.')
            return jsonify({'error': 'Supplier GTC file must be uploaded'}), 400

        if index_name:
            if abb_gtc:
                logger.error('Cannot upload ABB GTC file when index_name is selected.')
                return jsonify({'error': 'Cannot upload ABB GTC file when index_name is selected'}), 400
            files = [supplier_gtc]
            file_names = ['supplier_gtc']
        else:
            if not abb_gtc:
                logger.error('ABB GTC file must be uploaded if index_name is not provided.')
                return jsonify({'error': 'ABB GTC file must be uploaded if index_name is not provided'}), 400
            files = [abb_gtc, supplier_gtc]
            file_names = ['abb_gtc', 'supplier_gtc']

        if any(file.filename == '' for file in files):
            logger.error('One or both files are missing filenames.')
            return jsonify({'error': 'One or both files are missing filenames'}), 400

        global blob_list
        timestamp = datetime.now().strftime("%Y%m%d%H%M")
        folder_name = f"rfq_{timestamp}_{str(uuid.uuid4())}"

        file_info = []

        for file, name in zip(files, file_names):
            blob_name = f"{folder_name}/{file.filename}"
            blob_list.append(blob_name)
            logger.info(f'Blob Name {blob_name}')
            blobstorageprocessor.upload_blob(AZURE_BLOB_STORAGE_CONTAINER_NAME, blob_name, file)
            logger.info(f'File "{file.filename}" uploaded successfully.')

            file_url = f'{AZURE_BLOB_STORAGE_CONTAINER_URL}/{blob_name}'
            file_info.append({
                'name': name,
                'original_filename': file.filename,
                'blob_name': blob_name,
                'file_url': file_url
            })

        return jsonify({
            'success': 'Files uploaded successfully',
            'files': file_info
        }), 200

    except Exception as e:
        logger.error(f'Error in upload_files route: {str(e)}')
        return jsonify({'error': f'An error occurred: {str(e)}'}), 500

def process_entity_with_retries(entity, agent_additions, agent_others, tools, prompts, max_retries=1, backoff_factor=0.5):
    retries = 0
    while retries <= max_retries:
        try:
            # Select the appropriate agent based on the entity
            if entity == "ADDITIONS IN SUPPLIER GTC":
                # Use the template for "Additions in Supplier GTC"
                agentexecutor = AgentExecutor(agent=agent_additions, tools=tools, verbose=True, handle_parsing_errors=True)
                input_prompt = "Please analyze the documents and provide the additions."
            
            else:
                # Use the template for other entities
                agentexecutor = AgentExecutor(agent=agent_others, tools=tools, verbose=True, handle_parsing_errors=True)
                input_prompt = f"Conduct a comparative analysis of the {entity} clauses in both the ABB General Terms and Conditions (GTC) and the Supplier GTC. In your analysis, please address the following key aspects: {prompts[entity]}. Please ensure that your response is comprehensive, addressing each of these elements in detail."

            with get_openai_callback() as cb:
                response = agentexecutor.invoke({
                    "input": input_prompt
                })
                
                total_token = cb.total_tokens
                total_cost = cb.total_cost

            # Convert and return the response
            answer = response['output']
            
            if entity == "ADDITIONS IN SUPPLIER GTC":
                answer = string_to_dict(answer)
                # Append similarities, removals, and differences
                answer["similarities"] = []
                answer["removals"] = []
                answer["differences"] = []
            else:
                answer = convert_to_dict(answer)
                
            return entity, answer, total_token, total_cost
        
        except Exception as e:
            retries += 1
            if retries > max_retries:
                raise e
            time.sleep(backoff_factor * (2 ** retries))

    raise Exception(f"Max retries exceeded for entity: {entity}")

def is_effectively_empty(text):
    """
    Check if the text is effectively empty or contains only insignificant content.
    Returns a tuple (is_empty, patterns_found).
    """
    insignificant_patterns = [
        r'Sheet:\s*Sheet\d*',             # Excel sheet names like 'Sheet: Sheet1'
        r'Slide\s*\d+',                   # PowerPoint slide labels like 'Slide 1'
        r'Title Slide',                   # PowerPoint title slides
        r'Page\s*\d+',                    # Page numbers like 'Page 1'
        r'^\s*$',
        r'^$'                         # Empty lines or lines with only whitespace
        r'\n',                            # Newlines
        r'\r\n',                          # Carriage returns
        r'\t',                            # Tabs
        r'^Generated by.*$',              # Auto-generated messages
        r'^Created on.*$',                # Timestamps or creation info
        r'^Document\d*$',                 # Default Word document titles like 'Document1'
        r'^Untitled Document$',           # 'Untitled Document' in Word
        r'^Type text here$',              # Placeholder text in Word
        r'^[,\s]*$',                      # Lines with only commas or whitespace (CSV)
        r'^(?:\w+,?)*$',                  # Lines with only column headers (CSV)
        r'N/A',                           # Placeholder values
        r'No Data',                       # Placeholder indicating absence of data
        # Add more patterns as needed
    ]

    # Patterns found in the text
    patterns_found = []

    # Check for each pattern in the text
    for pattern in insignificant_patterns:
        matches = re.findall(pattern, text, flags=re.MULTILINE | re.IGNORECASE)
        if matches:
            patterns_found.extend(matches)

    # Remove insignificant content using regex substitution
    cleaned_text = re.sub('|'.join(insignificant_patterns), '', text, flags=re.MULTILINE | re.IGNORECASE)
    cleaned_text = cleaned_text.strip()

    # Set a threshold for minimal meaningful content length
    minimal_length = 10  # Adjust as needed

    # Determine if text is effectively empty
    is_empty = len(cleaned_text) < minimal_length

    return is_empty, patterns_found

def parse_blob_url(blob_url):
    parsed_url = urlparse(blob_url)
    path_parts = parsed_url.path.lstrip('/').split('/')
    container_name = path_parts[0]
    blob_name = '/'.join(path_parts[1:])
    return container_name, blob_name

@app.route('/process', methods=['GET', 'POST'])
def process_files():
    try:
        index_name = request.json.get('index_name')
        supplier_gtc_url = request.json.get('supplier_gtc_url')
        company_gtc_url = request.json.get('company_gtc_url')

        # Ensure the Supplier GTC URL is provided
        if not supplier_gtc_url:
            return jsonify({'error': 'Supplier GTC file URL must be provided'}), 400

        # Define files based on the provided inputs
        files = [{"name": "Supplier-GTC", "url": supplier_gtc_url}]
        if company_gtc_url:
            files.append({"name": "ABB-GTC", "url": company_gtc_url})
        elif index_name:
            files.append({"name": "ABB-GTC", "vector_store_index": index_name})

        # Load, split, and create retrievers for each document
        tools = []
        for file in files:
            if 'vector_store_index' in file:
                # Download FAISS index from Azure Blob Storage
                vector_download = blobstorageprocessor.download_blob(AZURE_VECTOR_STORAGE_CONTAINER_NAME, index_name)
                stream = vector_download.readall()
                retriever_store = FAISS.deserialize_from_bytes(embeddings=embeddings, serialized=stream)
                retriever = retriever_store.as_retriever(search_type="similarity", search_kwargs={"k": 10})
                blob_name = index_name + '.pdf'
                abb_entities = extract_entities_from_blob(AZURE_PRELOADED_CONTAINER_NAME, blob_name)
                abb_entities["24. ADDITIONS IN SUPPLIER GTC"] = []
                abb_entities = abb_entities.keys()
            else:
                # Process and create retrievers for new documents
                container_name, blob_name = parse_blob_url(file['url'])
                stream = blobstorageprocessor.download_blob(container_name, blob_name)
                if file['name'] == 'ABB-GTC':
                    blob_url = file['url']
                    abb_entities = extract_entities_from_blob(container_name, blob_name)
                    abb_entities["24. ADDITIONS IN SUPPLIER GTC"] = []
                    abb_entities = abb_entities.keys()

                if file['url'].endswith(('.pdf', '.docx')):
                    # Load and split PDF or DOCX using appropriate methods
                    try:
                        blob_url = blobstorageprocessor.get_blob_sas_url(container_name, blob_name)
                        result = documentintelligence.analyze_document_word(blob_url)
                        text = result.content
                    except Exception as e:
                        return jsonify({'error': f'Error reading file: {str(e)}'}), 500

                # Handle other file formats (e.g., .doc, .xls, .xlsx, etc.)
                elif file['url'].endswith('.doc'):
                    # Create a temporary directory
                    logger.info('DOC File is being Processed')
                    blob_url = blobstorageprocessor.get_blob_sas_url(AZURE_BLOB_STORAGE_CONTAINER_NAME, blob_name)
                    response = requests.get(blob_url)

                    # Use io.BytesIO instead of writing to a local file
                    doc_content = io.BytesIO(response.content)
                    # logger.info("The content of this document is:",response.content)
                    # Create a temporary file for the conversion process
                    with NamedTemporaryFile(suffix='.doc', delete=False) as temp_doc:
                        temp_doc.write(doc_content.getvalue())
                        temp_doc_path = temp_doc.name

                    # Convert DOC to DOCX using LibreOffice
                    output_dir = os.path.dirname(temp_doc_path)
                    subprocess.call(['lowriter', '--headless', '--convert-to', 'docx', '--outdir', output_dir, temp_doc_path])


                    # The converted file will have the same name but with .docx extension
                    logger.info("document converted to doc:")
                    temp_docx_path = temp_doc_path.replace('.doc', '.docx')

                    # Read the DOCX file
                    text = docx2txt2.extract_text(temp_docx_path)
                    logger.info("The content of this document is:"+text)

                    # Extract text from the DOCX
                    # doc_text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])

                    os.remove(temp_doc_path)
                    os.remove(temp_docx_path)
                elif file['url'].endswith('.xls'):
                    logger.info("XLS Processing Start")
                    logger.info("-" * 30)
                    excel_content = BytesIO(stream.readall())

                    workbook = xlrd.open_workbook(file_contents=excel_content.getvalue())
                    all_sheets_content = ""

                    for sheet_index in range(workbook.nsheets):
                        sheet = workbook.sheet_by_index(sheet_index)
                        sheet_content = f"Sheet: {sheet.name}\n"

                        for row in range(sheet.nrows):
                            row_values = []
                            for col in range(sheet.ncols):
                                cell_value = sheet.cell_value(row, col)
                                if sheet.cell_type(row, col) == xlrd.XL_CELL_DATE:
                                    cell_value = xlrd.xldate.xldate_as_datetime(cell_value, workbook.datemode).strftime('%Y-%m-%d %H:%M:%S')
                                row_values.append(str(cell_value))
                            sheet_content += "\t".join(row_values) + "\n"

                        all_sheets_content += sheet_content + "\n\n"

                    text = all_sheets_content
                    logger.info("XLS Processing End")
                    logger.info("-" * 30)
                elif file['url'].endswith('.xlsx'):
                    logger.info("XLSX Processing Start")
                    logger.info("-" * 30)
                    excel_content = BytesIO(stream.readall())
                    excel_df = pd.read_excel(excel_content, sheet_name=None)
                    all_sheets_content = ""
                    for sheet_name, df in excel_df.items():
                        all_sheets_content += f"Sheet: {sheet_name}\n"
                        all_sheets_content += df.to_csv(index=False)
                    text = all_sheets_content
                    logger.info("XLSX Processing End")
                    logger.info("-" * 30)
                elif file['url'].endswith('.csv'):
                    logger.info("CSV Processing Start")
                    csv_content = BytesIO(stream.readall())
                    logger.info("CSV contents read")
                    csv_df = pd.read_csv(csv_content)
                    text = csv_df.to_csv(index=False)
                    logger.info("CSV Processing End")
                elif file['url'].endswith('.pptx'):
                    logger.info("PPTX Processing Start")
 
                    # Fetch the blob stream
                    logger.info("PPTX file read from Blob URL")
                    pptx_content = BytesIO(stream.readall())
                    logger.info("PPTX contents read")
 
                    # Load the presentation
                    presentation = Presentation(pptx_content)
 
                    # Extract text from slides
                    text = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
 
                    logger.info("PPTX Processing End")

                else:
                    return jsonify({'error': 'Unsupported file format'}), 400
                # After extracting text for each file
                if file['name'] == 'Supplier-GTC':
                    supplier_text = text  # Save supplier text for later use

                # Check if the text is effectively empty
                is_empty, patterns_found = is_effectively_empty(text)
                if is_empty:
                    # Handle the empty text case
                    logger.info(f"The text for {file['name']} is empty or insignificant.")
                    
                    # Create an EmptyRetriever
                    retriever = EmptyRetriever()
                else:
                # Split the extracted text into chunks for further analysis
                    headers_to_split_on = [
                        ("#", "Header 1"),
                        ("##", "Header 2"),
                        ("###", "Header 3"),
                    ]

                    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                    splits = markdown_splitter.split_text(text)

                    # Create FAISS index and retrievers
                    faiss_index = FAISS.from_documents(splits, embedding=embeddings)
                    retriever = faiss_index.as_retriever()
                    bm25_retriever = BM25Retriever.from_texts([text])
                    retriever = EnsembleRetriever(retrievers=[bm25_retriever, retriever], weights=[0.5, 0.5])# [Your existing code for handling different file types]
                print(f"Retriever created for {file}")
                # ...
                tools.append(
                    Tool(
                        args_schema=DocumentInput,
                        name=file["name"],
                        description=f"useful when you want to answer questions about {file['name']}",
                        func=RetrievalQA.from_chain_type(llm=llm, retriever=retriever),
                    )
                )
        # Now, check if Supplier-GTC text is effectively empty
        is_empty, patterns_found = is_effectively_empty(supplier_text)
        patterns_str = ', '.join(set(patterns_found))  # Convert patterns to a string

        # Define the first system template for "Additions in Supplier GTC"
        template_additions = """You are an AI assistant specializing in document comparison.

**Task**:

- Compare two documents: **ABB-GTC** and **Supplier-GTC**.
- Identify clauses present in **Supplier-GTC** but **absent** in **ABB-GTC** (Additions).

**Instructions**:

- Provide concise summaries for each addition.
- Do **not** include any extra text or explanations.
- Return the result as a list.
- Ensure all strings are properly quoted.
- **Escape any curly braces in your output.**

**Output Format**:

[
    "Summary of addition 1,
    Summary of addition 2"
  ]

"""



        # Define the second system template for other entities
        template_others = """You are an AI assistant specializing in document comparison.

**Task**:

- Compare two documents: **ABB-GTC** and **Supplier-GTC**.
- Categorize clauses into:
  - **Similarities**: Summary of clauses that are common to both documents.
  - **Additions**: Summary of clauses present in **Supplier-GTC** but **absent** in **ABB-GTC**.
  - **Removals**: Summary of clauses present in **ABB-GTC** but **absent** in **Supplier-GTC**.
  - **Differences**: Summary of clauses that are common to both but have differences.

**Instructions**:

- Return the result as a valid JSON object.
- Do **not** include any extra text or explanations.
- Do **not** include code snippets or function calls in your response.
- Ensure all strings are properly quoted.
- **Escape any curly braces in your output.**
- Highlight important keywords or terms using double asterisks (**), e.g., **Delivery**, **Payment Terms**.

**Output Format**:

```json
{{
  "similarities": [
    "Summary of similarity 1"
  ],
  "additions": [
    "Summary of addition 1"
  ],
  "removals": [
    "Summary of removal 1"
  ],
  "differences": [
    "Summary of difference 1"
  ]
}}
"""
        if is_empty:
            # Modify the template_others to include the instruction
            template_others_modified = """You are an AI assistant specializing in document comparison.

**Task**:

- The **Supplier-GTC** text contains the following insignificant patterns or is empty: {{patterns_str}}.
- Move all relevant **ABB-GTC** clauses to the "Removals" section.
- Leave all other sections empty.

**Instructions**:

- Return the result as a valid JSON object.
- Do **not** include any extra text or explanations.
- Do **not** include code snippets or function calls in your response.
- Ensure all strings are properly quoted.
- **Escape any curly braces in your output.**

**Output Format**:

```json
{{
  "similarities": [],
  "additions": [],
  "removals": [
    "Summary of removal 1"
  ],
  "differences": []
}}
"""
        else:
            # Use the original template_others
            template_others_modified = template_others
        # Create two prompt templates
        prompt_template_additions = ChatPromptTemplate.from_messages([
            ("system", template_additions),
            ("human", "{input}"),
            MessagesPlaceholder(variable_name="agent_scratchpad")
        ])

        prompt_template_others = ChatPromptTemplate.from_messages([
            ("system", template_others_modified),
            ("human", "{input}"),
            MessagesPlaceholder(variable_name="agent_scratchpad")
        ])

        # Create two agents with their respective prompt templates
        agent_additions = create_openai_functions_agent(
            tools=tools,
            llm=llm,
            prompt=prompt_template_additions
        )

        agent_others = create_openai_functions_agent(
            tools=tools,
            llm=llm,
            prompt=prompt_template_others
        )

        # Extract entities
        entities = [re.split(r"^\d+.", entity)[1].strip() for entity in abb_entities]

        # Initialize variables for storing the results
        result_dict = {}
        total_cost = 0
        total_token = 0
        max_workers = 8  # Adjust based on your system capacity

        # Execute the agent for each entity concurrently
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_entity = {
                executor.submit(process_entity_with_retries, entity,agent_additions,agent_others, tools, prompts): entity for entity in entities
            }
            for future in as_completed(future_to_entity):
                entity = future_to_entity[future]
                try:
                    entity, result, tokens, cost = future.result()
                    result_dict[entity] = result
                    total_token += tokens
                    total_cost += cost
                except Exception as e:
                    result_dict[entity] = {'error': str(e)}

        sorted_dict = {key: result_dict[key] for key in entities if key in result_dict}
        df = pd.DataFrame(sorted_dict)
        df = df.T
        df = df.reset_index()
        df.rename(columns={"index": "Entities"}, inplace=True)
        if 'error' in df.columns:
            df.drop(['error'], axis=1, inplace=True)
        lang_selected = 'en'
        output = create_excel_with_formatting_local(df, lang_selected, sheet_name='Output')
        unique_id = str(uuid.uuid4())[:8]
        blob_name = f"Comparison_{unique_id}.xlsx"
        blob_service_client = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONNECTION_STRING_DOWNLOAD_EXCEL)
        container_client = blob_service_client.get_container_client(AZURE_DOWNLOAD_EXCEL_STORAGE_CONTAINER_NAME)
        blob_client = container_client.get_blob_client(blob_name)
        blob_client.upload_blob(output, overwrite=True)

        if(total_cost <= 0 and total_token > 0):
            total_cost = 1
            
        return jsonify({'result': sorted_dict, 'total_cost': total_cost, 'total_token': total_token, 'blob_name': blob_name}), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/material_process', methods=['POST'])
def material_process():
    try:
        choice = request.form.get('choice')
        if not choice:
            return jsonify({"error": "Format type not specified"}), 400

        if choice == 'with_format':
            if 'layout_file' not in request.files:
                return jsonify({"error": "Layout file is missing"}), 400

            layout_file = request.files['layout_file']
            if layout_file.filename == '':
                return jsonify({"error": "No layout file selected"}), 400

            try:
                layout_df = pd.read_excel(layout_file)

            except Exception as e:
                return jsonify({"error": f"Failed to read layout file: {str(e)}"}), 400

            selected_columns = request.form.getlist('selected_columns')
            if not selected_columns:
                return jsonify({"error": "No columns selected for comparison"}), 400

            comparison_files = request.files.getlist('comparison_files')
            if not comparison_files:
                return jsonify({"error": "No comparison files uploaded"}), 400

            file_column_selections = {}
            for file in comparison_files:
                columns = request.form.getlist(f'file_columns_{file.filename}')
                if not columns:
                    return jsonify({"error": f"Invalid or missing file format. Please upload a file in the correct format and try again."}), 400
                file_column_selections[file.filename] = columns

            return material_process_files(layout_df, comparison_files, selected_columns, file_column_selections)

        elif choice == 'without_format':
            uploaded_files = request.files.getlist('comparison_files')
            if len(uploaded_files) < 2:
                return jsonify({"error": "Please upload at least two files for comparison"}), 400

            file_column_selections = {}
            for file in uploaded_files:
                columns = request.form.getlist(f'file_columns_{file.filename}')
                if not columns:
                    return jsonify({"error": f"Invalid or missing file format. Please upload a file in the correct format and try again."}), 400
                file_column_selections[file.filename] = columns

            return materials_process_files_without_format(uploaded_files, file_column_selections)
        else:
            return jsonify({"error": f"Invalid choice: {choice}"}), 400

    except Exception as e:
        return jsonify({"error": f"An unexpected error occurred: {str(e)}"}), 500

def material_process_files(layout_df, comparison_files, selected_columns, file_column_selections):
    try:
        comparison_dfs = {}
        all_matches = {}
        
        for file in comparison_files:
            try:
                comparison_df = read_excel_preserve_format(file)
                column_mapping, unmapped_columns = matcher.map_columns(layout_df, comparison_df, selected_columns)
                matches = matcher.find_row_matches(layout_df, comparison_df, column_mapping, selected_columns)
                
                comparison_dfs[file.filename] = (comparison_df, unmapped_columns, column_mapping)
                all_matches[file.filename] = matches
            except Exception as e:
                return jsonify({"error": f"Error processing file {file.filename}: {str(e)}"}), 400

        combined_df = matcher.create_combined_dataframe(layout_df, comparison_dfs, all_matches, selected_columns, file_column_selections)
        
        # Reorder columns
        reordered_columns = selected_columns.copy()
        for file_name in comparison_dfs.keys():
            reordered_columns.extend([f"{file_name}_{col}" for col in file_column_selections[file_name]])

        # Create an OrderedDict to maintain column order
        result_data = OrderedDict()
        for column in combined_df.columns:
            result_data[column] = combined_df[column].tolist()

        return app.response_class(
            response=json.dumps({"result": result_data}, cls=OrderedEncoder),
            status=200,
            mimetype='application/json'
        )

    except Exception as e:
        return jsonify({"error": f"Error in processing files: {str(e)}"}), 500

def materials_process_files_without_format(uploaded_files, file_column_selections):
    try:
        if len(uploaded_files) < 2:
            return jsonify({"error": "Please upload at least two files for comparison"}), 400

        try:
            reference_df = read_excel_preserve_format(uploaded_files[0])
        except Exception as e:
            return jsonify({"error": f"Failed to read reference file: {str(e)}"}), 400

        comparison_files = uploaded_files[1:]

        target_columns = ["Description", "Qty", "UOM", "Manufacturer"]
        reference_mapping = matcher.find_similar_columns(reference_df, target_columns)

        comparison_dfs = {}
        all_matches = {}

        for file in comparison_files:
            try:
                comparison_df = pd.read_excel(file)
                column_mapping, _ = matcher.map_columns(reference_df, comparison_df, list(reference_mapping.values()))
                
                if not column_mapping:
                    return jsonify({"error": f"Could not map columns for file: {file.filename}"}), 400
                
                matches = matcher.find_row_matches(reference_df, comparison_df, reference_mapping, column_mapping)
                comparison_dfs[file.filename] = (comparison_df, list(set(comparison_df.columns) - set(column_mapping.values())), column_mapping)
                all_matches[file.filename] = matches
            except Exception as e:
                return jsonify({"error": f"Error processing file {file.filename}: {str(e)}"}), 400

        # Create combined DataFrame
        combined_df = reference_df[list(reference_mapping.values())].copy()

        # Add columns from all files, including the first one
        for file_name, selected_cols in file_column_selections.items():
            if not selected_cols:
                return jsonify({"error": f"No columns selected for file: {file_name}"}), 400
            for col in selected_cols:
                combined_df[f"{file_name}_{col}"] = None

        # Populate data for the first file (reference file)
        first_file_name = uploaded_files[0].filename
        try:
            for col in file_column_selections[first_file_name]:
                if col not in reference_df.columns:
                    return jsonify({"error": f"Column '{col}' not found in reference file"}), 400
                combined_df[f"{first_file_name}_{col}"] = reference_df[col]
        except KeyError as e:
            return jsonify({"error": f"Error accessing column in reference file: {str(e)}"}), 400

        # Populate data for comparison files
        for file_name, (comparison_df, _, column_mapping) in comparison_dfs.items():
            matches = all_matches[file_name]
            try:
                for i, row in combined_df.iterrows():
                    match = next((m for m in matches if m[0] == i), None)
                    if match and match[1] is not None:
                        comparison_row = comparison_df.iloc[match[1]]
                        for col in file_column_selections[file_name]:
                            mapped_col = column_mapping.get(col, col)
                            if mapped_col not in comparison_row.index:
                                return jsonify({"error": f"Column '{mapped_col}' not found in file: {file_name}"}), 400
                            combined_df.at[i, f"{file_name}_{col}"] = comparison_row.get(mapped_col, None)
            except Exception as e:
                return jsonify({"error": f"Error populating data from file {file_name}: {str(e)}"}), 400

        # Reorder columns
        try:
            reordered_columns = list(reference_mapping.values())
            for file_name in [first_file_name] + list(comparison_dfs.keys()):
                reordered_columns.extend([f"{file_name}_{col}" for col in file_column_selections[file_name]])
            
            combined_df = combined_df[reordered_columns]
        except KeyError as e:
            return jsonify({"error": f"Error reordering columns: {str(e)}"}), 400

        # Convert DataFrame to OrderedDict
        try:
            result_data = OrderedDict()
            for column in combined_df.columns:
                result_data[column] = combined_df[column].tolist()

            return app.response_class(
                response=json.dumps({"result": result_data}, cls=OrderedEncoder),
                status=200,
                mimetype='application/json'
            )
        except Exception as e:
            return jsonify({"error": f"Error preparing final data: {str(e)}"}), 500

    except Exception as e:
        return jsonify({"error": f"An unexpected error occurred: {str(e)}"}), 500
    
@app.route('/get_columns', methods=['POST'])
def get_columns():
    file = request.files['file']
    df = pd.read_excel(file)
    columns = df.columns.tolist()
    return jsonify(columns)

@app.route('/find_similar_columns', methods=['POST'])
def api_find_similar_columns():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files['file']
    target_columns = ["Description", "Qty", "UOM", "Manufacturer"]
    
    if not target_columns:
        return jsonify({"error": "No target columns provided"}), 400
    try:
        df = pd.read_excel(file)
        column_mapping = matcher.find_similar_columns(df, target_columns)
        if len(column_mapping) == 0:
            return jsonify({"error": "No matching columns found"}), 400
        else:
            return jsonify({"column_mapping": column_mapping})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/facilities_process', methods=['POST'])
def process():
    if 'format_file' not in request.files:
        return jsonify({"error": "No format file provided"}), 400
    
    if 'pdf_files' not in request.files:
        return jsonify({"error": "No PDF files provided"}), 400
    
    format_file = request.files['format_file']
    pdf_files = request.files.getlist('pdf_files')
    
    tools = []
    failed_files = []
    choice = request.form.get('choice')            
    
    if choice == "format2":
        try:
            # Initialize matcher
            matcher = facilitesMatcher()
            
            # Read layout file
            layout_df = pd.read_excel(format_file)
            selected_columns = layout_df.columns.tolist()
            
            # Process comparison files
            comparison_dfs = {}
            file_column_selections = {}
            
            for file in pdf_files:
                comparison_df = pd.read_excel(file)
                
                # Handle format if needed
                for col in selected_columns:
                    if any(comparison_df.iloc[:, 0] == col):
                        first_match_row = comparison_df.index[comparison_df.iloc[:, 0] == col][0]
                        comparison_df = comparison_df.iloc[first_match_row + 1:].set_axis(comparison_df.iloc[first_match_row], axis=1).reset_index(drop=True)
                        break
                
                comparison_df = fix_duplicate_columns(comparison_df)
                
                # Find similar columns
                similar_columns = set()
                for selected_column in selected_columns:
                    similar_cols = matcher.find_similar_columns(comparison_df, [selected_column])
                    similar_columns.update(similar_cols.values())
                
                additional_columns = [col for col in comparison_df.columns if col not in similar_columns]
                file_column_selections[file.filename] = additional_columns
                comparison_dfs[file.filename] = comparison_df
            
            # Process matches
            all_matches = {}
            processed_comparison_dfs = {}
            
            for file_name, comparison_df in comparison_dfs.items():
                column_mapping, unmapped_columns = matcher.map_columns(layout_df, comparison_df, selected_columns)
                matches = matcher.find_row_matches(layout_df, comparison_df, column_mapping, selected_columns)
                processed_comparison_dfs[file_name] = (comparison_df, unmapped_columns, column_mapping)
                all_matches[file_name] = matches
            
            # Create combined DataFrame
            combined_df = matcher.create_combined_dataframe(
                layout_df, 
                processed_comparison_dfs, 
                all_matches, 
                selected_columns, 
                file_column_selections
            )
            
            # Convert DataFrame to JSON
            return jsonify({
                "success": True,
                "data": combined_df.to_dict(orient='records')
            })
            
        except Exception as e:
            logging.error("Error processing files: %s", e)
            return jsonify({
                "error": str(e)
            }), 500
        
    result_facilites = read_xls_to_prompts(format_file)
    prompts = result_facilites[0]
    non_filled = result_facilites[1]
    error = result_facilites[2]
    
    logger.info(f"non_filled: {non_filled}")
    logger.info(f"Prompts: {prompts}")
    if error:
        return jsonify({"error": "Provide proper excel format"}), 400
  
    for file in pdf_files:
        tool = process_single_file(file)
        if tool:
            tools.append(tool)
        else:
            failed_files.append(file.filename)
    
    if not tools:
        return jsonify({
            "error": "Failed to process all files",
            "failed_files": failed_files
        }), 400
    print("prompt dict : ", prompts)
    prompts_dict = {str(idx+1): prompt for idx, prompt in enumerate(prompts)}
    prompts_json_str = json.dumps({"prompts": prompts_dict}, indent=2)
    
    # Create dynamic JSON structure based on non_filled list
    dynamic_structure = ", ".join(f'"{item}": "..."' for item in non_filled) if non_filled else '"amount": "...", "quantity": "...", "total": "..."'
    question = prompts_json_str
    
    prompt_template = f"""
    You are a precise JSON generator. Your task is to extract information from the document as per the provided prompts, and calculate a Total for each prompt.

    Rules:
    1. Always respond with ONLY a valid JSON object, no additional text
    2. For each prompt, extract values based on the provided fields
    3. If information is not found, use ""
    4. total = amount * quantity
    4. Remove currency symbols before multiplication, then reapply to total

    Prompts: {question}

    Respond with ONLY this JSON structure:
    {{
      "prompts": {{
        1: {{
            {dynamic_structure},
            total: "..."
        }},
        2: {{
            {dynamic_structure},
            total: "..."
        }},
        ... (continue for all prompts)
      }}
    }}
    """
    
    logger.info(f"Prompt Template: {prompt_template}")
    
    query = prompt_template
    
    async def process_tool(tool, query):
        try:
            response = tool.func.invoke({"query": query})
            return tool.name, response["result"]
        except Exception as e:
            logging.error(f"Error processing tool {tool.name}: {str(e)}")
            return tool.name, None
    
    async def process_all_tools(tools, query):
        tasks = [process_tool(tool, query) for tool in tools]
        return await asyncio.gather(*tasks)
    
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    results = loop.run_until_complete(process_all_tools(tools, query))

    final_response = {
        "prompt_details": {str(i+1): prompt for i, prompt in enumerate(prompts)},
        "results": {}
    }
    processing_errors = []
    parsing_errors = {}
    
    for file_name, result in results:
        if result is None:
            processing_errors.append(file_name)
            continue
        
        try:
            parsed_result = extract_and_parse_json(result)
            if parsed_result and 'prompts' in parsed_result and parsed_result['prompts']:
                final_response["results"][file_name] = parsed_result
            else:
                logging.error(f"Failed to extract valid JSON structure from {file_name}")
                processing_errors.append(file_name)
                parsing_errors[file_name] = {
                    "error": "Failed to extract valid JSON structure",
                    "raw_response": result[:200] + "..." if len(result) > 200 else result
                }
        except Exception as e:
            logging.error(f"Error processing result from {file_name}: {str(e)}")
            processing_errors.append(file_name)
            parsing_errors[file_name] = {
                "error": str(e),
                "raw_response": result[:200] + "..." if len(result) > 200 else result
            }
    
    if failed_files or processing_errors:
        final_response["errors"] = {
            "failed_to_process": failed_files,
            "failed_to_extract": processing_errors
        }
        if parsing_errors:
            final_response["errors"]["parsing_errors"] = parsing_errors

    return jsonify(final_response)


if __name__ == '__main__':
    app.run()